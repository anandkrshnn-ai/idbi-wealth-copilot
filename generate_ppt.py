import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Define colors
COLOR_DARK_BLUE = RGBColor(10, 37, 64)      # #0A2540
COLOR_ACCENT_ORANGE = RGBColor(255, 107, 0)  # #FF6B00
COLOR_TEAL = RGBColor(0, 128, 128)          # #008080
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_BG = RGBColor(244, 246, 248)    # #F4F6F8
COLOR_TEXT_DARK = RGBColor(33, 37, 41)       # #212529
COLOR_TEXT_MUTED = RGBColor(108, 117, 125)   # #6C757D
COLOR_BORDER = RGBColor(222, 226, 230)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="AI WEALTH COPILOT PROTOTYPE"):
    # Header bar background shape
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_DARK_BLUE
    header_box.line.color.rgb = COLOR_DARK_BLUE

    # Category label
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT_ORANGE
    p_cat.font.name = 'Arial'

    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE
    p_title.font.name = 'Arial'

def add_card(slide, left, top, width, height, bg_color=COLOR_WHITE, border_color=COLOR_BORDER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_bullet_point(tf, bold_prefix, text_body, size=14, space_after=10):
    p = tf.add_paragraph()
    p.space_after = Pt(space_after)
    
    run1 = p.add_run()
    run1.text = bold_prefix
    run1.font.bold = True
    run1.font.size = Pt(size)
    run1.font.color.rgb = COLOR_TEXT_DARK
    run1.font.name = 'Arial'
    
    run2 = p.add_run()
    run2.text = text_body
    run2.font.bold = False
    run2.font.size = Pt(size)
    run2.font.color.rgb = COLOR_TEXT_DARK
    run2.font.name = 'Arial'

def main():
    prs = Presentation()
    # Widescreen format (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # -------------------------------------------------------------
    # SLIDE 1: Title Slide & Problem Statement
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BLUE)
    
    # Large Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "AI WEALTH COPILOT PROTOTYPE"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.font.name = 'Arial'
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Conversational, Behavior-Aware Compliance Advisory"
    p1_sub.font.size = Pt(20)
    p1_sub.font.color.rgb = COLOR_ACCENT_ORANGE
    p1_sub.font.name = 'Arial'
    p1_sub.space_before = Pt(10)
    
    # Meta / Team Details
    meta_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(5.5), Inches(3.8))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    
    p = tf_meta.paragraphs[0]
    p.text = "TEAM DETAILS"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    
    p_t = tf_meta.add_paragraph()
    p_t.text = (
        "• Team Leader: Anandakrishnan Damodaran\n"
        "• Problem Statement: Digital Wealth Management\n"
        "• Organisation: Independent Builder\n"
        "• Location: Chennai, Tamil Nadu"
    )
    p_t.font.size = Pt(14)
    p_t.font.color.rgb = COLOR_WHITE
    p_t.space_before = Pt(6)
    p_t.line_spacing = 1.3
    
    p_links = tf_meta.add_paragraph()
    p_links.text = "\nONLINE LINKS\n• Repo: github.com/anandkrshnn-ai/ai-wealth-copilot-hackathon-prototype\n• Live App: anandkrshnn-ai.github.io/ai-wealth-copilot-hackathon-prototype/"
    p_links.font.size = Pt(12)
    p_links.font.color.rgb = COLOR_WHITE
    p_links.space_before = Pt(12)
    
    # Problem Statement Card (Right Side)
    add_card(slide1, Inches(6.8), Inches(2.6), Inches(5.7), Inches(4.0), bg_color=RGBColor(20, 50, 80), border_color=None)
    prob_box = slide1.shapes.add_textbox(Inches(7.1), Inches(2.8), Inches(5.1), Inches(3.6))
    tf_prob = prob_box.text_frame
    tf_prob.word_wrap = True
    
    p_prob_hdr = tf_prob.paragraphs[0]
    p_prob_hdr.text = "THE PROBLEM"
    p_prob_hdr.font.bold = True
    p_prob_hdr.font.size = Pt(14)
    p_prob_hdr.font.color.rgb = COLOR_ACCENT_ORANGE
    
    p_prob_body = tf_prob.add_paragraph()
    p_prob_body.text = (
        "Banks want to offer digital wealth management across their retail customer base to increase "
        "onboarding funnel conversions. However, they face high regulatory compliance risks regarding "
        "suitability (SEBI/RBI rules) and cannot afford unauthorized automated advice.\n\n"
        "Rigid rule engines frustrate customers, while black-box generative AI tools hallucinate recommendations. "
        "A hybrid model resolves this tension by using AI for natural intent capture and audit rationales, "
        "with a rules-first compliance engine checking every checkout."
    )
    p_prob_body.font.size = Pt(12)
    p_prob_body.font.color.rgb = COLOR_WHITE
    p_prob_body.space_before = Pt(10)
    p_prob_body.line_spacing = 1.2

    # Disclaimer on Title Slide
    disc_box1 = slide1.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    tf_disc1 = disc_box1.text_frame
    tf_disc1.word_wrap = True
    p_disc1 = tf_disc1.paragraphs[0]
    p_disc1.text = "This is an independent prototype created for demonstration purposes. It is not affiliated with or endorsed by any bank."
    p_disc1.font.size = Pt(10)
    p_disc1.font.italic = True
    p_disc1.font.color.rgb = COLOR_BORDER
    p_disc1.font.name = 'Arial'
    
    # -------------------------------------------------------------
    # SLIDE 2: Brief About the Idea
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_header(slide2, "Brief About the Idea - Digital Advisory Re-imagined")
    
    # Left Card: Project Scope
    add_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    left_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.6))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p = tf_left.paragraphs[0]
    p.text = "AI WEALTH COPILOT PROTOTYPE"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_DARK_BLUE
    p.space_after = Pt(12)
    
    add_bullet_point(tf_left, "Guardrailed Digital Wealth Advisor: ", "Embedded directly inside mobile banking journeys, converting free-form goals into suitability-compliant outcomes.")
    add_bullet_point(tf_left, "AI + Rules Hybrid Engine: ", "AI manages multilingual onboarding conversation, while a hardcoded rules engine holds final veto authority.")
    add_bullet_point(tf_left, "RM Integration Gateway: ", "Regulated, complex, or high-value cases automatically hand off to relationship managers with full context summaries.")
    
    # Right Card: The Control Interface
    add_card(slide2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0), bg_color=COLOR_DARK_BLUE, border_color=None)
    right_box = slide2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.6))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p = tf_right.paragraphs[0]
    p.text = "THE COMPLIANCE ADVANTAGE"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p.space_after = Pt(12)
    
    p_rd1 = tf_right.add_paragraph()
    p_rd1.text = "1. Safe Onboarding Conversions"
    p_rd1.font.bold = True
    p_rd1.font.size = Pt(13)
    p_rd1.font.color.rgb = COLOR_WHITE
    
    p_rd2 = tf_right.add_paragraph()
    p_rd2.text = "Uses Gemini 1.5 Flash (with Bedrock Claude 3.5 Sonnet proxy pathways) to extract intent structure and explain rationales, keeping friction low."
    p_rd2.font.size = Pt(12)
    p_rd2.font.color.rgb = COLOR_LIGHT_BG
    p_rd2.space_after = Pt(14)
    
    p_rd3 = tf_right.add_paragraph()
    p_rd3.text = "2. Enforcing Strict Suitability Policies"
    p_rd3.font.bold = True
    p_rd3.font.size = Pt(13)
    p_rd3.font.color.rgb = COLOR_WHITE
    
    p_rd4 = tf_right.add_paragraph()
    p_rd4.text = "Enforces regulatory rules (SEBI/RBI/IRDA) to prevent unauthorized advice, log compliance overrides, and maintain a full audit trail."
    p_rd4.font.size = Pt(12)
    p_rd4.font.color.rgb = COLOR_LIGHT_BG
    
    # -------------------------------------------------------------
    # SLIDE 3: Opportunities
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_header(slide3, "Opportunities & Unique Selling Propositions")
    
    # 3 Column Bento Box layout
    add_card(slide3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8))
    box1 = slide3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.2), Inches(4.4))
    tf_b1 = box1.text_frame
    tf_b1.word_wrap = True
    tf_b1.paragraphs[0].text = "HOW IS IT DIFFERENT?"
    tf_b1.paragraphs[0].font.bold = True
    tf_b1.paragraphs[0].font.color.rgb = COLOR_ACCENT_ORANGE
    add_bullet_point(tf_b1, "Workflow Copilot: ", "Not a generic LLM chatbot — it has explicit, hardcoded suitability rules.", size=13)
    add_bullet_point(tf_b1, "Behavioral Inputs: ", "Uses goals + affordability + transaction logs, not just static questionnaires.", size=13)

    add_card(slide3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
    box2 = slide3.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.4))
    tf_b2 = box2.text_frame
    tf_b2.word_wrap = True
    tf_b2.paragraphs[0].text = "HOW WILL IT SOLVE IT?"
    tf_b2.paragraphs[0].font.bold = True
    tf_b2.paragraphs[0].font.color.rgb = COLOR_TEAL
    add_bullet_point(tf_b2, "Scale Advisory: ", "Brings personalized advisory to mass-retail segments without linear RM headcount growth.", size=13)
    add_bullet_point(tf_b2, "Better RM Leads: ", "Pre-screens customer risk, suitability, and readiness to invest before RM handoff.", size=13)

    add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8))
    box3 = slide3.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.2), Inches(4.4))
    tf_b3 = box3.text_frame
    tf_b3.word_wrap = True
    tf_b3.paragraphs[0].text = "UNIQUE SELLING POINTS"
    tf_b3.paragraphs[0].font.bold = True
    tf_b3.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE
    add_bullet_point(tf_b3, "Guardrailed GenAI: ", "Every advice recommendation must pass suitability parameters before the user sees it.", size=13)
    add_bullet_point(tf_b3, "Explainability: ", "Delivers explicit reason cards (Goal fit, Risk fit, Product fit) for RMs and auditors.", size=13)

    # -------------------------------------------------------------
    # SLIDE 4: List of Features Offered
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_header(slide4, "List of Features Offered by the Solution")
    
    # 2x4 Grid Layout of cards (8 features)
    features = [
        ("Goal-based Onboarding", "Capture goals, liquidity horizons, and timelines conversing in native language."),
        ("Risk Profiling", "Combine KYC, demographics, income, and transaction behavior to assign bands."),
        ("Product Suitability Engine", "Map mutual funds, bonds, deposits, and insurance safely to risk bands."),
        ("Multilingual UX", "Supports English, Hindi, Marathi, and Gujarati conversational chat interfaces."),
        ("RM Handoff Gateway", "Auto-flags complex scenarios (like PMS >50L) and forwards logs to RMs."),
        ("Explainability Cards", "Delivers customer-safe and compliance-specific rationales for every advice."),
        ("Audit Console Log", "Exposes AI prompts, parameters, rules overrides, and API traces to auditors."),
        ("Behavioral scoring checks", "Scans account logs to penalize cheque bounces and reward savings habits.")
    ]
    
    width_box = Inches(2.7)
    height_box = Inches(2.2)
    positions = [
        (Inches(0.6), Inches(1.8)), (Inches(3.6), Inches(1.8)), (Inches(6.6), Inches(1.8)), (Inches(9.6), Inches(1.8)),
        (Inches(0.6), Inches(4.4)), (Inches(3.6), Inches(4.4)), (Inches(6.6), Inches(4.4)), (Inches(9.6), Inches(4.4))
    ]
    
    for i, (title, desc) in enumerate(features):
        pos = positions[i]
        add_card(slide4, pos[0], pos[1], width_box, height_box)
        tb = slide4.shapes.add_textbox(pos[0] + Inches(0.15), pos[1] + Inches(0.15), width_box - Inches(0.3), height_box - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_BLUE
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 5: Process Flow Diagram
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_LIGHT_BG)
    add_header(slide5, "Process Flow Diagram - Advisory Onboarding Journey")
    
    # 6 blocks horizontally
    steps = [
        ("1. Consent & Log", "User logs in, opts into Wealth Copilot."),
        ("2. Goal Capture", "AI collects goals, horizons in natural language."),
        ("3. Suitability Engine", "Checks risk bands, credit scores, transaction history."),
        ("4. AI Rationales", "Generates explainability cards showing fits & rationales."),
        ("5. Guardrail Check", "Triggers direct checkout OR automatic RM escalation."),
        ("6. Order Executed", "Completes the transaction, schedules portfolio monitoring.")
    ]
    
    step_width = Inches(1.8)
    step_gap = Inches(0.2)
    left_start = Inches(0.6)
    top_pos = Inches(2.5)
    
    for i, (title, desc) in enumerate(steps):
        l = left_start + i * (step_width + step_gap)
        add_card(slide5, l, top_pos, step_width, Inches(3.2), bg_color=COLOR_WHITE)
        
        num_box = slide5.shapes.add_textbox(l + Inches(0.1), top_pos + Inches(0.15), Inches(1.6), Inches(0.4))
        p_num = num_box.text_frame.paragraphs[0]
        p_num.text = f"STAGE {i+1}"
        p_num.font.bold = True
        p_num.font.size = Pt(10)
        p_num.font.color.rgb = COLOR_ACCENT_ORANGE
        
        det_box = slide5.shapes.add_textbox(l + Inches(0.1), top_pos + Inches(0.6), Inches(1.6), Inches(2.4))
        tf = det_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK_BLUE
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_TEXT_DARK
        p_d.space_before = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 6: Screens / Journey View
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_LIGHT_BG)
    add_header(slide6, "Screens / Journey View - Customer Advisory Journey")
    
    # 4 columns for Onboarding, Profiling, Recommendation, Escalation
    journey_steps = [
        ("1. Onboarding", 
         "• Customer opens mobile app & opts in\n"
         "• Consent captured for profile-based advice\n"
         "• Starts conversational multilingual journey"),
        ("2. Profiling", 
         "• Collects goal horizon, budget, & timeline\n"
         "• Integrates dynamic account behaviour checks\n"
         "• Suitability logic segments risk bands"),
        ("3. Recommendation", 
         "• Presents explainable recommendations\n"
         "• Explains Goal, Risk, & Product fits\n"
         "• Policy guardrails verify eligible options"),
        ("4. Escalation / Execution", 
         "• Auto-escapes complex or regulated cases\n"
         "• RMs receive complete context summaries\n"
         "• Executes & sets up event-driven nudges")
    ]
    
    w_j = Inches(2.7)
    g_j = Inches(0.2)
    start_j = Inches(0.6)
    
    for i, (title, points) in enumerate(journey_steps):
        l = start_j + i * (w_j + g_j)
        add_card(slide6, l, Inches(1.8), w_j, Inches(4.2))
        
        tb = slide6.shapes.add_textbox(l + Inches(0.15), Inches(2.0), w_j - Inches(0.3), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_DARK_BLUE
        p.space_after = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = points
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_DARK
        p_desc.line_spacing = 1.3

    footer_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8))
    p_foot = footer_box.text_frame.paragraphs[0]
    p_foot.text = "A guided, compliant journey from discovery to execution — with AI for scale and RM oversight for trust."
    p_foot.font.italic = True
    p_foot.font.size = Pt(13)
    p_foot.font.color.rgb = COLOR_TEXT_MUTED
    p_foot.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 7: Architecture Diagram
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_LIGHT_BG)
    add_header(slide7, "Architecture Diagram of the Proposed Solution")
    
    # Let's describe the architecture layers beautifully in columns
    add_card(slide7, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.8))
    tf_arch1 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3.2), Inches(4.4)).text_frame
    tf_arch1.word_wrap = True
    tf_arch1.paragraphs[0].text = "CHANNEL LAYER"
    tf_arch1.paragraphs[0].font.bold = True
    tf_arch1.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE
    tf_arch1.paragraphs[0].space_after = Pt(10)
    add_bullet_point(tf_arch1, "Client Mobile App: ", "Host client chat.", size=12)
    add_bullet_point(tf_arch1, "RM Branch Portal: ", "Receives escalations.", size=12)
    add_bullet_point(tf_arch1, "API Gateway: ", "Routes secure client sessions.", size=12)

    add_card(slide7, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.6), bg_color=COLOR_DARK_BLUE, border_color=None)
    tf_arch2 = slide7.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.2), Inches(4.2)).text_frame
    tf_arch2.word_wrap = True
    tf_arch2.paragraphs[0].text = "AI & LOGIC LAYER"
    tf_arch2.paragraphs[0].font.bold = True
    tf_arch2.paragraphs[0].font.color.rgb = COLOR_ACCENT_ORANGE
    tf_arch2.paragraphs[0].space_after = Pt(10)
    
    def add_bullet_white_arch(tf, prefix, desc):
        p = tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = prefix
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_LIGHT_BG

    add_bullet_white_arch(tf_arch2, "Orchestrator: ", "Session state & chat routing.")
    add_bullet_white_arch(tf_arch2, "Policy Engine: ", "Hardcoded SEBI suitability & risk check rules.")
    add_bullet_white_arch(tf_arch2, "Bedrock / Gemini: ", "For intent extraction & explainability.")

    add_card(slide7, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.8))
    tf_arch3 = slide7.shapes.add_textbox(Inches(9.0), Inches(1.8), Inches(3.2), Inches(4.4)).text_frame
    tf_arch3.word_wrap = True
    tf_arch3.paragraphs[0].text = "DATA & CORE LAYER"
    tf_arch3.paragraphs[0].font.bold = True
    tf_arch3.paragraphs[0].font.color.rgb = COLOR_TEAL
    tf_arch3.paragraphs[0].space_after = Pt(10)
    add_bullet_point(tf_arch3, "Profile Master: ", "Holds customer demographic info.", size=12)
    add_bullet_point(tf_arch3, "Core Banking (CBS): ", "Ingests active balances and transaction warning logs.", size=12)
    add_bullet_point(tf_arch3, "RDS Database: ", "Logs audit trails securely.", size=12)

    note_box = slide7.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6))
    p_note = note_box.text_frame.paragraphs[0]
    p_note.text = "Note: The architecture isolates bank data inside private subnets, routing to Amazon Bedrock via secure VPC endpoints."
    p_note.font.italic = True
    p_note.font.size = Pt(12)
    p_note.font.color.rgb = COLOR_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 8: Technologies Used
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_LIGHT_BG)
    add_header(slide8, "Technologies to be Used in the Solution")
    
    # 4 Column Stack Comparison
    tech_data = [
        ("FRONTEND & UX", "• HTML5 & CSS3\n• Vanilla Javascript\n• Mobile banking integration\n• Multilingual UX templates"),
        ("AI & REASONING", "• Gemini 1.5 Flash NLU\n• Amazon Bedrock Claude 3.5\n• Structured JSON Output\n• Guardrails & RAG"),
        ("BACKEND & DATA", "• API Gateway & BFF services\n• Customer profile masters\n• PostgreSQL database logs\n• EventBus notifications"),
        ("SECURITY & COMP.", "• OAuth & SSO auth layers\n• SSL & KMS envelope encryption\n• DPDP Act consent tracking\n• Auditor audit console logging")
    ]
    
    w_card = Inches(2.7)
    g_card = Inches(0.2)
    start_left = Inches(0.6)
    
    for i, (title, list_items) in enumerate(tech_data):
        l = start_left + i * (w_card + g_card)
        add_card(slide8, l, Inches(1.8), w_card, Inches(4.8))
        
        tb = slide8.shapes.add_textbox(l + Inches(0.15), Inches(2.0), w_card - Inches(0.3), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_DARK_BLUE
        p.space_after = Pt(12)
        
        p_list = tf.add_paragraph()
        p_list.text = list_items
        p_list.font.size = Pt(12)
        p_list.font.color.rgb = COLOR_TEXT_DARK
        p_list.line_spacing = 1.3

    # -------------------------------------------------------------
    # SLIDE 9: Performance, Cost & Benchmarks
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_LIGHT_BG)
    add_header(slide9, "Prototype Benchmarking, Costs & Future Plans")
    
    # Left Card: Benchmarks & Specs
    add_card(slide9, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    perf_l = slide9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.6))
    tf_pl = perf_l.text_frame
    tf_pl.word_wrap = True
    tf_pl.paragraphs[0].text = "BENCHMARKS & PERFORMANCE REPORT"
    tf_pl.paragraphs[0].font.bold = True
    tf_pl.paragraphs[0].font.size = Pt(14)
    tf_pl.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE
    tf_pl.paragraphs[0].space_after = Pt(10)
    add_bullet_point(tf_pl, "Fast Response: ", "Takes < 5 seconds for generating advisory cards.", size=12)
    add_bullet_point(tf_pl, "100% Guardrail: ", "Policy rules intercept vulnerability segments and check surplus affordability.", size=12)
    add_bullet_point(tf_pl, "Explainability: ", "Delivers 3 critical reason cards for goal, risk, and product alignment.", size=12)
    add_bullet_point(tf_pl, "1-Click Handoff: ", "RMs receive context payloads with complete scoring audit logs.", size=12)

    # Right Card: Costs & Roadmaps
    add_card(slide9, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.0), bg_color=COLOR_DARK_BLUE, border_color=None)
    perf_r = slide9.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.6))
    tf_pr = perf_r.text_frame
    tf_pr.word_wrap = True
    tf_pr.paragraphs[0].text = "IMPLEMENTATION COST & ROADMAP"
    tf_pr.paragraphs[0].font.bold = True
    tf_pr.paragraphs[0].font.size = Pt(14)
    tf_pr.paragraphs[0].font.color.rgb = COLOR_ACCENT_ORANGE
    tf_pr.paragraphs[0].space_after = Pt(10)
    
    def add_bullet_white_p8(tf, prefix, desc):
        p = tf.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = prefix
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_WHITE
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_LIGHT_BG

    add_bullet_white_p8(tf_pr, "Development: ", "6-8 weeks from prototype sandbox verification to pilot release.")
    add_bullet_white_p8(tf_pr, "Low Cost: ", "Uses existing cloud environments, caching data during test phases.")
    add_bullet_white_p8(tf_pr, "Near-term: ", "Mutual Fund Master integrations, portfolio rebalancing triggers.")
    add_bullet_white_p8(tf_pr, "Long-term: ", "Audit logs monitoring dashboard, fairness checks, structured investments.")

    # -------------------------------------------------------------
    # SLIDE 10: Snapshots of the Prototype
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10, COLOR_LIGHT_BG)
    add_header(slide10, "Snapshots of the Active Advisory Prototype")
    
    # 4 placeholder cards side by side where the user can paste/overlay screenshots of the app
    screenshot_slots = [
        ("Goal & Profile Discovery", "Client goal text box & language selector panel"),
        ("Suitability Analysis", "Calculated score indicators & behavioral score meters"),
        ("Advisory Rationale", "Customer-safe & SEBI compliance explainability cards"),
        ("RM Review & Escalation", "Connect-to-RM gateway and secure handoff payloads")
    ]
    
    w_s = Inches(2.7)
    g_s = Inches(0.2)
    start_s = Inches(0.6)
    
    for i, (title, desc) in enumerate(screenshot_slots):
        l = start_s + i * (w_s + g_s)
        add_card(slide10, l, Inches(1.8), w_s, Inches(4.5), bg_color=COLOR_WHITE, border_color=COLOR_ACCENT_ORANGE)
        
        tb = slide10.shapes.add_textbox(l + Inches(0.15), Inches(2.0), w_s - Inches(0.3), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Snapshot 0{i+1}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_ORANGE
        
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.bold = True
        p_title.font.size = Pt(13)
        p_title.font.color.rgb = COLOR_DARK_BLUE
        p_title.space_before = Pt(6)
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"\n[Insert Screenshot Here]\n\n{desc}"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = COLOR_TEXT_MUTED
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.space_before = Pt(20)

    # -------------------------------------------------------------
    # SLIDE 11: Slide Links & Project Assets
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11, COLOR_DARK_BLUE)
    
    end_box = slide11.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.2))
    tf11 = end_box.text_frame
    tf11.word_wrap = True
    p11 = tf11.paragraphs[0]
    p11.text = "SUBMISSION LINKS & ASSETS"
    p11.font.size = Pt(36)
    p11.font.bold = True
    p11.font.color.rgb = COLOR_WHITE
    p11.font.name = 'Arial'
    
    # 3 Large Cards for links
    add_card(slide11, Inches(0.8), Inches(2.8), Inches(3.6), Inches(3.8), bg_color=RGBColor(20, 50, 80), border_color=None)
    c1 = slide11.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(3.2), Inches(3.4))
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "GITHUB PUBLIC REPO"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p_desc = tf_c1.add_paragraph()
    p_desc.text = "\nSource code, compliance tests, and target specs:\n\nhttps://github.com/anandkrshnn-ai/ai-wealth-copilot-hackathon-prototype"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = COLOR_WHITE
    
    add_card(slide11, Inches(4.8), Inches(2.8), Inches(3.6), Inches(3.8), bg_color=RGBColor(20, 50, 80), border_color=None)
    c2 = slide11.shapes.add_textbox(Inches(5.0), Inches(3.0), Inches(3.2), Inches(3.4))
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "LIVE DEPLOYMENT LINK"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEAL
    p_desc = tf_c2.add_paragraph()
    p_desc.text = "\nInteract with the live static mobile UI and Auditor console:\n\nhttps://anandkrshnn-ai.github.io/ai-wealth-copilot-hackathon-prototype/"
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = COLOR_WHITE
    
    add_card(slide11, Inches(8.8), Inches(2.8), Inches(3.6), Inches(3.8), bg_color=RGBColor(20, 50, 80), border_color=None)
    c3 = slide11.shapes.add_textbox(Inches(9.0), Inches(3.0), Inches(3.2), Inches(3.4))
    tf_c3 = c3.text_frame
    tf_c3.word_wrap = True
    p = tf_c3.paragraphs[0]
    p.text = "DEMO VIDEO LINK (3 MIN)"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_ACCENT_ORANGE
    p_desc = tf_c3.add_paragraph()
    p_desc.text = "\nWatch the live demo recording showing NLU goal parser, transaction scoring rules, and RM escalation guardrails."
    p_desc.font.size = Pt(12)
    p_desc.font.color.rgb = COLOR_WHITE
    
    # Disclaimer on Final Slide
    disc_box11 = slide11.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    tf_disc11 = disc_box11.text_frame
    tf_disc11.word_wrap = True
    p_disc11 = tf_disc11.paragraphs[0]
    p_disc11.text = "This is an independent prototype created for demonstration purposes. It is not affiliated with or endorsed by any bank."
    p_disc11.font.size = Pt(10)
    p_disc11.font.italic = True
    p_disc11.font.color.rgb = COLOR_BORDER
    p_disc11.font.name = 'Arial'

    # Save presentation
    filepath = "AI_Wealth_Copilot_Pitch_Deck.pptx"
    prs.save(filepath)
    print(f"Success! PowerPoint presentation saved at: {os.path.abspath(filepath)}")

if __name__ == "__main__":
    main()
